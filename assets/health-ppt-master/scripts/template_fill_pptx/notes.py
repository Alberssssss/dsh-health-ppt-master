"""Speaker-notes parts for cloned slides.

Builds native PowerPoint notes-slide XML and the slide<->notesSlide<->notesMaster
relationships from a plan's ``notes`` field, reusing the SVG pipeline's notes
renderer so embedded notes also feed ``notes_to_audio.py``.
"""

from __future__ import annotations

from copy import deepcopy
import io
import posixpath
import re
import zipfile
from xml.etree import ElementTree as ET

from pptx import Presentation

from svg_to_pptx.pptx_notes import create_notes_slide_xml, markdown_to_plain_text

from .ooxml import (
    CT_NS,
    NS,
    NOTES_MASTER_CONTENT_TYPE,
    NOTES_MASTER_REL_TYPE,
    NOTES_SLIDE_CONTENT_TYPE,
    NOTES_SLIDE_REL_TYPE,
    REL_NS,
    SLIDE_REL_TYPE,
    THEME_CONTENT_TYPE,
    THEME_REL_TYPE,
    _qn,
    _xml_bytes,
)
from .package import (
    _add_content_type_override,
    _empty_relationships_root,
    _max_numeric_rid,
)


def _find_notes_master_target(entries: dict[str, bytes]) -> str | None:
    for name, data in entries.items():
        if not name.startswith("ppt/notesSlides/_rels/notesSlide") or not name.endswith(".xml.rels"):
            continue
        try:
            root = ET.fromstring(data)
        except ET.ParseError:
            continue
        for rel in root.findall(_qn(REL_NS, "Relationship")):
            if rel.attrib.get("Type") == NOTES_MASTER_REL_TYPE:
                return rel.attrib.get("Target")

    presentation_rels = entries.get("ppt/_rels/presentation.xml.rels")
    if not presentation_rels:
        return None
    try:
        root = ET.fromstring(presentation_rels)
    except ET.ParseError:
        return None
    for rel in root.findall(_qn(REL_NS, "Relationship")):
        if rel.attrib.get("Type") != NOTES_MASTER_REL_TYPE:
            continue
        target = rel.attrib.get("Target")
        if not target:
            return None
        if target.startswith("/"):
            target = target.lstrip("/")
        else:
            target = posixpath.normpath(posixpath.join("ppt", target))
        return posixpath.relpath(target, "ppt/notesSlides")
    return None


def _next_part_number(entries: dict[str, bytes], pattern: str) -> int:
    compiled = re.compile(pattern)
    return max(
        (int(match.group(1)) for name in entries if (match := compiled.match(name))),
        default=0,
    ) + 1


def _default_notes_master_parts() -> tuple[bytes, bytes, bytes]:
    """Return a standards-complete default notes master, rels, and theme.

    Delegate the verbose OOXML boilerplate to python-pptx so this path stays in
    lockstep with the main SVG exporter and does not maintain a hand-written
    notes-master template.
    """
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _ = slide.notes_slide
    payload = io.BytesIO()
    prs.save(payload)
    payload.seek(0)
    with zipfile.ZipFile(payload) as donor:
        return (
            donor.read("ppt/notesMasters/notesMaster1.xml"),
            donor.read("ppt/notesMasters/_rels/notesMaster1.xml.rels"),
            donor.read("ppt/theme/theme2.xml"),
        )


def _ensure_notes_master(
    entries: dict[str, bytes],
    presentation_root: ET.Element,
    presentation_rels_root: ET.Element,
    content_root: ET.Element,
) -> str:
    """Ensure a PowerPoint-standard notesMaster package and presentation edge."""
    master_xml, master_rels_xml, theme_xml = _default_notes_master_parts()
    default_master_root = ET.fromstring(master_xml)
    default_master_rels_root = ET.fromstring(master_rels_xml)
    existing = _find_notes_master_target(entries)

    if existing:
        master_part = posixpath.normpath(
            posixpath.join("ppt/notesSlides", existing)
        )
        if master_part not in entries:
            raise RuntimeError(
                f"Existing notesMaster relationship targets a missing part: {master_part}"
            )
        master_name = posixpath.basename(master_part)
        master_rels_part = posixpath.join(
            posixpath.dirname(master_part), "_rels", f"{master_name}.rels"
        )
        master_root = ET.fromstring(entries[master_part])
        if master_root.find("p:notesStyle", NS) is None:
            default_notes_style = default_master_root.find("p:notesStyle", NS)
            if default_notes_style is None:
                raise RuntimeError("Default notes master has no p:notesStyle")
            master_root.append(deepcopy(default_notes_style))
            entries[master_part] = _xml_bytes(master_root)
        master_rels_root = (
            ET.fromstring(entries[master_rels_part])
            if master_rels_part in entries
            else deepcopy(default_master_rels_root)
        )
    else:
        master_num = _next_part_number(
            entries, r"^ppt/notesMasters/notesMaster(\d+)\.xml$"
        )
        master_part = f"ppt/notesMasters/notesMaster{master_num}.xml"
        master_rels_part = (
            f"ppt/notesMasters/_rels/notesMaster{master_num}.xml.rels"
        )
        entries[master_part] = master_xml
        master_rels_root = deepcopy(default_master_rels_root)

    theme_rels = [
        rel
        for rel in master_rels_root.findall(_qn(REL_NS, "Relationship"))
        if rel.attrib.get("Type") == THEME_REL_TYPE
    ]
    if len(theme_rels) > 1:
        raise RuntimeError("notesMaster must not have multiple theme relationships")
    if not theme_rels:
        theme_rel = ET.SubElement(
            master_rels_root,
            _qn(REL_NS, "Relationship"),
            {
                "Id": f"rId{_max_numeric_rid(master_rels_root) + 1}",
                "Type": THEME_REL_TYPE,
                "Target": "",
            },
        )
    else:
        theme_rel = theme_rels[0]

    raw_theme_target = theme_rel.attrib.get("Target", "")
    if raw_theme_target:
        theme_part = posixpath.normpath(
            posixpath.join(posixpath.dirname(master_part), raw_theme_target)
        )
    else:
        theme_num = _next_part_number(entries, r"^ppt/theme/theme(\d+)\.xml$")
        theme_part = f"ppt/theme/theme{theme_num}.xml"
        theme_rel.set(
            "Target",
            posixpath.relpath(theme_part, posixpath.dirname(master_part)),
        )
    if theme_part not in entries:
        entries[theme_part] = theme_xml

    entries[master_rels_part] = _xml_bytes(master_rels_root)
    _add_content_type_override(content_root, master_part, NOTES_MASTER_CONTENT_TYPE)
    _add_content_type_override(content_root, theme_part, THEME_CONTENT_TYPE)

    presentation_notes_rels = [
        rel
        for rel in presentation_rels_root.findall(_qn(REL_NS, "Relationship"))
        if rel.attrib.get("Type") == NOTES_MASTER_REL_TYPE
    ]
    if presentation_notes_rels:
        master_rel = presentation_notes_rels[0]
        for duplicate in presentation_notes_rels[1:]:
            presentation_rels_root.remove(duplicate)
    else:
        master_rel = ET.SubElement(
            presentation_rels_root,
            _qn(REL_NS, "Relationship"),
            {"Id": f"rId{_max_numeric_rid(presentation_rels_root) + 1}"},
        )
    master_rel.set("Type", NOTES_MASTER_REL_TYPE)
    master_rel.set("Target", posixpath.relpath(master_part, "ppt"))
    master_rel_id = master_rel.attrib["Id"]

    notes_master_id_lst = presentation_root.find("p:notesMasterIdLst", NS)
    if notes_master_id_lst is None:
        notes_master_id_lst = ET.Element(_qn(NS["p"], "notesMasterIdLst"))
        children = list(presentation_root)
        insert_at = next(
            (
                index
                for index, child in enumerate(children)
                if child.tag in {
                    _qn(NS["p"], "handoutMasterIdLst"),
                    _qn(NS["p"], "sldIdLst"),
                }
            ),
            len(children),
        )
        presentation_root.insert(insert_at, notes_master_id_lst)
    else:
        for child in list(notes_master_id_lst):
            notes_master_id_lst.remove(child)
    ET.SubElement(
        notes_master_id_lst,
        _qn(NS["p"], "notesMasterId"),
        {_qn(NS["r"], "id"): master_rel_id},
    )

    return posixpath.relpath(master_part, "ppt/notesSlides")


def _create_notes_rels_xml(slide_number: int, notes_master_target: str | None) -> bytes:
    root = _empty_relationships_root()
    if notes_master_target:
        ET.SubElement(
            root,
            _qn(REL_NS, "Relationship"),
            {
                "Id": "rId1",
                "Type": NOTES_MASTER_REL_TYPE,
                "Target": notes_master_target,
            },
        )
        slide_rel_id = "rId2"
    else:
        slide_rel_id = "rId1"
    ET.SubElement(
        root,
        _qn(REL_NS, "Relationship"),
        {
            "Id": slide_rel_id,
            "Type": SLIDE_REL_TYPE,
            "Target": f"../slides/slide{slide_number}.xml",
        },
    )
    return _xml_bytes(root)


def _verify_standard_notes_package(
    entries: dict[str, bytes],
    presentation_root: ET.Element,
    presentation_rels_root: ET.Element,
    content_root: ET.Element,
) -> None:
    """Reject notes packages that strict PowerPoint/iOS readers cannot consume."""
    note_parts = sorted(
        name
        for name in entries
        if re.fullmatch(r"ppt/notesSlides/notesSlide\d+\.xml", name)
    )
    if not note_parts:
        return

    problems: list[str] = []
    master_rels = [
        rel
        for rel in presentation_rels_root.findall(_qn(REL_NS, "Relationship"))
        if rel.attrib.get("Type") == NOTES_MASTER_REL_TYPE
    ]
    master_rel_id = master_rels[0].attrib.get("Id", "") if len(master_rels) == 1 else ""
    master_target = master_rels[0].attrib.get("Target", "") if len(master_rels) == 1 else ""
    if len(master_rels) != 1 or not master_rel_id or not master_target:
        problems.append("presentation must have exactly one notesMaster relationship")
    master_part = (
        master_target.lstrip("/")
        if master_target.startswith("/")
        else posixpath.normpath(posixpath.join("ppt", master_target))
    )

    notes_ids = presentation_root.findall(
        "./p:notesMasterIdLst/p:notesMasterId", NS
    )
    notes_id_rids = [node.attrib.get(_qn(NS["r"], "id"), "") for node in notes_ids]
    if len(notes_ids) != 1 or notes_id_rids != [master_rel_id]:
        problems.append("notesMasterIdLst does not reference the notesMaster relationship")

    master_data = entries.get(master_part)
    if not master_data:
        problems.append(f"missing notesMaster part: {master_part}")
    else:
        master_root = ET.fromstring(master_data)
        if master_root.find("p:notesStyle", NS) is None:
            problems.append(f"{master_part} has no p:notesStyle")

    master_rels_part = posixpath.join(
        posixpath.dirname(master_part),
        "_rels",
        f"{posixpath.basename(master_part)}.rels",
    )
    master_rels_data = entries.get(master_rels_part)
    if not master_rels_data:
        problems.append(f"missing notesMaster relationships: {master_rels_part}")
    else:
        master_rels_root = ET.fromstring(master_rels_data)
        theme_rels = [
            rel
            for rel in master_rels_root.findall(_qn(REL_NS, "Relationship"))
            if rel.attrib.get("Type") == THEME_REL_TYPE
        ]
        if len(theme_rels) != 1:
            problems.append("notesMaster must have exactly one theme relationship")
        elif theme_rels[0].attrib.get("Target"):
            theme_part = posixpath.normpath(
                posixpath.join(
                    posixpath.dirname(master_part),
                    theme_rels[0].attrib["Target"],
                )
            )
            if theme_part not in entries:
                problems.append(f"missing notesMaster theme: {theme_part}")

    overrides = {
        node.attrib.get("PartName", "").lstrip("/"): node.attrib.get("ContentType", "")
        for node in content_root.findall(_qn(CT_NS, "Override"))
    }
    if overrides.get(master_part) != NOTES_MASTER_CONTENT_TYPE:
        problems.append(f"invalid Content Type for {master_part}")
    invalid_rels_overrides = sorted(
        name
        for name in overrides
        if name.endswith(".rels")
        and name.startswith(("ppt/notesSlides/", "ppt/notesMasters/"))
    )
    if invalid_rels_overrides:
        problems.append(
            "notes relationship parts must use the rels Default, not Overrides: "
            + ", ".join(invalid_rels_overrides)
        )

    for note_part in note_parts:
        note_root = ET.fromstring(entries[note_part])
        placeholder_types = {
            node.attrib.get("type", "")
            for node in note_root.findall(".//p:ph", NS)
        }
        missing = {"sldImg", "body", "sldNum"} - placeholder_types
        if missing:
            problems.append(
                f"{note_part} is missing placeholders: {', '.join(sorted(missing))}"
            )
        if note_root.find("p:clrMapOvr", NS) is None:
            problems.append(f"{note_part} has no p:clrMapOvr")
        if overrides.get(note_part) != NOTES_SLIDE_CONTENT_TYPE:
            problems.append(f"invalid Content Type for {note_part}")

        note_rels_part = posixpath.join(
            posixpath.dirname(note_part),
            "_rels",
            f"{posixpath.basename(note_part)}.rels",
        )
        note_rels_data = entries.get(note_rels_part)
        if not note_rels_data:
            problems.append(f"missing notesSlide relationships: {note_rels_part}")
            continue
        note_rels_root = ET.fromstring(note_rels_data)
        note_master_links = [
            rel
            for rel in note_rels_root.findall(_qn(REL_NS, "Relationship"))
            if rel.attrib.get("Type") == NOTES_MASTER_REL_TYPE
        ]
        slide_links = [
            rel
            for rel in note_rels_root.findall(_qn(REL_NS, "Relationship"))
            if rel.attrib.get("Type") == SLIDE_REL_TYPE
        ]
        if len(note_master_links) != 1:
            problems.append(f"{note_rels_part} must have one notesMaster relationship")
        if len(slide_links) != 1:
            problems.append(f"{note_rels_part} must have one slide relationship")

    if problems:
        details = "\n".join(f"  - {problem}" for problem in problems)
        raise RuntimeError(
            "PPTX package contains a non-standard speaker-notes structure; "
            "iOS may render the deck as blank:\n" + details
        )


def _slide_rels_with_notes(
    rels_bytes: bytes | None,
    *,
    slide_number: int,
    notes_text: str,
    notes_master_target: str | None,
) -> tuple[bytes, dict[str, bytes]]:
    root = ET.fromstring(rels_bytes) if rels_bytes else _empty_relationships_root()
    for rel in list(root.findall(_qn(REL_NS, "Relationship"))):
        if rel.attrib.get("Type") == NOTES_SLIDE_REL_TYPE:
            root.remove(rel)

    note_entries: dict[str, bytes] = {}
    notes_text = notes_text.strip()
    if notes_text:
        rel_id = f"rId{_max_numeric_rid(root) + 1}"
        notes_part = f"ppt/notesSlides/notesSlide{slide_number}.xml"
        notes_rels_part = f"ppt/notesSlides/_rels/notesSlide{slide_number}.xml.rels"
        ET.SubElement(
            root,
            _qn(REL_NS, "Relationship"),
            {
                "Id": rel_id,
                "Type": NOTES_SLIDE_REL_TYPE,
                "Target": f"../notesSlides/notesSlide{slide_number}.xml",
            },
        )
        plain_notes = markdown_to_plain_text(notes_text)
        note_entries[notes_part] = create_notes_slide_xml(slide_number, plain_notes).encode("utf-8")
        note_entries[notes_rels_part] = _create_notes_rels_xml(slide_number, notes_master_target)

    return _xml_bytes(root), note_entries
